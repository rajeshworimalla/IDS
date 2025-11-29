#!/usr/bin/env python3
"""
Script to create PowerPoint presentation for IDS Capstone Project
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(0, 51, 102)  # Dark blue
    bullet_color = RGBColor(0, 0, 0)   # Black
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "An Intrusion Detection System (IDS)"
    subtitle.text = "by [Team Member 1], [Team Member 2], [Team Member 3]"
    
    # Slide 2: Requirements - Why was this IDS created?
    slide = add_text_slide(prs, "Requirements", 
        "Why was this IDS created?\n\n" +
        "With the increasing frequency and sophistication of cyber attacks, organizations need " +
        "intelligent systems to monitor and protect their networks in real-time. Traditional " +
        "security measures are often reactive, generating too many false positives and " +
        "overwhelming security teams with alert fatigue.\n\n" +
        "We wanted to create a comprehensive intrusion detection system that combines " +
        "rule-based detection with machine learning to accurately identify and respond to " +
        "network threats. The purpose of this system is to provide real-time network " +
        "monitoring, automated threat detection, and instant alerts for security teams.\n\n" +
        "As a group, we wanted to tackle the challenges of network security, machine learning " +
        "integration, and real-time system development.")
    
    # Slide 3: Perspective
    slide = add_bullet_slide(prs, "Perspective", [
        "IDS is an intrusion detection system designed to monitor network traffic and detect security threats in real-time",
        "Priority is on detecting malicious activity including DoS attacks, port scans, brute force attempts, and unauthorized access",
        "Should be a web application with a backend database where security administrators can monitor network activity, view alerts, and manage blocked IPs"
    ])
    
    # Slide 4: User Characteristics
    slide = add_bullet_slide(prs, "User Characteristics", [
        "Ability to create an account; working email",
        "Two categories of users:\n  • Security Administrators (full access to all features)\n  • Security Analysts (monitoring and viewing capabilities)",
        "Language options only in English",
        "System should be a web application so it is accessible from various locations, including remote monitoring",
        "System needs to securely and reliably handle sensitive network data including packet captures, IP addresses, attack patterns, and security logs",
        "The system should use a database to store packet data, attack records, and system configuration",
        "System should be able to capture and analyze network packets in real-time without impacting network performance",
        "System should have an intuitive, easy to use user interface that displays network statistics, attack alerts, packet streams, and blocked IPs",
        "System should be implemented with the latest security measures including input sanitization, JWT authentication, encrypted connections (TLS), and secure API endpoints"
    ])
    
    # Slide 5: User Requirements
    slide = add_bullet_slide(prs, "User Requirements - Functional Requirements", [
        "Real-Time Packet Capture: System must capture network packets from the network interface in real-time",
        "Attack Detection: System must detect 6 types of attacks:\n  • Denial of Service (DoS) attacks\n  • Port scanning and reconnaissance (Probe)\n  • Remote to Local (R2L) unauthorized access\n  • User to Root (U2R) privilege escalation\n  • Brute force login attempts\n  • Unknown/unclassified malicious activity",
        "Machine Learning Integration: System must use ML models to classify packets as malicious or benign",
        "Rule-Based Detection: System must use rule-based detectors to identify specific attack patterns",
        "Real-Time Alerts: System must send instant alerts via WebSocket when attacks are detected",
        "IP Blocking: System must automatically block malicious IP addresses using firewall rules",
        "Dashboard Visualization: System must display network statistics, attack trends, and real-time packet streams",
        "User Authentication: System must support user registration, login, and session management",
        "Historical Data: System must store and allow querying of historical packet and attack data",
        "IP Management: System must allow administrators to manually block/unblock IP addresses"
    ])
    
    # Slide 6: System Architecture - Technology Stack
    slide = add_bullet_slide(prs, "System Architecture - Technology Stack", [
        "Frontend:\n  • React 19 with TypeScript\n  • Material-UI for components\n  • Recharts for data visualization\n  • Socket.io Client for real-time updates",
        "Backend:\n  • Node.js with TypeScript\n  • Express.js REST API\n  • Socket.io WebSocket server\n  • Mongoose for MongoDB\n  • libpcap for packet capture",
        "Machine Learning:\n  • Python 3.8+ with Flask\n  • scikit-learn for ML models\n  • Rule-based attack detectors",
        "Infrastructure:\n  • MongoDB for data storage\n  • Redis for caching\n  • iptables/ipset for firewall blocking"
    ])
    
    # Slide 7: Architecture Diagram
    slide = add_text_slide(prs, "System Architecture Diagram",
        "Frontend (React)\n" +
        "  ├─ Dashboard\n" +
        "  ├─ Monitoring\n" +
        "  ├─ Activities\n" +
        "  └─ Blocker\n\n" +
        "        ↓ WebSocket / REST API\n\n" +
        "Backend (Node.js/TypeScript)\n" +
        "  ├─ Packet Capture\n" +
        "  ├─ Aggregator\n" +
        "  └─ Blocker\n\n" +
        "        ↓\n\n" +
        "Prediction Service (Python/Flask)\n" +
        "  ├─ Rule-Based Detectors\n" +
        "  ├─ ML Models (Binary + Multiclass)\n" +
        "  └─ Attack Classifier\n\n" +
        "        ↓\n\n" +
        "Data Storage Layer\n" +
        "  ├─ MongoDB (Packets)\n" +
        "  ├─ Redis (Cache)\n" +
        "  └─ Firewall Rules")
    
    # Slide 8: Key Features - Attack Detection
    slide = add_bullet_slide(prs, "Key Features - Attack Detection", [
        "Dual Detection System: Combines rule-based pattern matching with machine learning models for accurate threat identification",
        "6 Attack Types: Detects DoS, Probe, R2L, U2R, Brute Force, and Unknown attacks",
        "Real-Time Processing: Analyzes packets as they are captured with minimal latency",
        "Confidence Scoring: ML models provide confidence scores for each detection"
    ])
    
    # Slide 9: Key Features - Automated Response & UI
    slide = add_bullet_slide(prs, "Key Features - Automated Response & User Interface", [
        "Automated Response:\n  • Automatic IP Blocking: Malicious IPs are automatically blocked via firewall rules\n  • Configurable Thresholds: Administrators can adjust detection sensitivity\n  • Whitelist Support: Trusted IPs can be whitelisted to prevent false positives",
        "User Interface:\n  • Real-Time Dashboard: Live network statistics and attack alerts\n  • Packet Monitoring: Stream of captured packets with filtering capabilities\n  • Historical Analysis: View past attacks and network activity\n  • IP Management: Manual blocking/unblocking of IP addresses\n  • Alert System: Instant popup notifications for detected attacks"
    ])
    
    # Slide 10: Machine Learning Models
    slide = add_bullet_slide(prs, "Machine Learning Models", [
        "Binary Classification Model:\n  • Purpose: Distinguish malicious vs benign traffic\n  • Accuracy: 100% on test data\n  • Algorithm: Random Forest Classifier\n  • Features: 41 network features extracted from packets",
        "Multiclass Classification Model:\n  • Purpose: Classify specific attack types\n  • Classes: Normal, DoS, Probe, R2L, U2R, Brute Force\n  • Algorithm: Random Forest Classifier\n  • Integration: Works alongside rule-based detectors",
        "Rule-Based Detectors:\n  • Port Scan Detector: Tracks unique ports per source IP\n  • DoS Detector: Monitors connection rates and flood patterns\n  • Brute Force Detector: Tracks failed login attempts\n  • R2L Detector: Identifies unauthorized access attempts\n  • U2R Detector: Detects privilege escalation patterns"
    ])
    
    # Slide 11: Security Measures
    slide = add_bullet_slide(prs, "Security Measures", [
        "Input Sanitization: All user inputs are validated and sanitized",
        "JWT Authentication: Secure token-based authentication system",
        "TLS Encryption: Encrypted connections for data in transit (production ready)",
        "Password Hashing: bcrypt for secure password storage",
        "Rate Limiting: API endpoints protected against brute force",
        "Secure API Design: RESTful API with proper error handling",
        "Firewall Integration: Direct integration with system firewall for IP blocking"
    ])
    
    # Slide 12: System Workflow
    slide = add_bullet_slide(prs, "System Workflow", [
        "Packet Capture: Network packets are captured from the network interface using libpcap",
        "Feature Extraction: Packet features are extracted (IPs, ports, protocols, timing)",
        "Detection Analysis:\n  • Rule-based detectors analyze patterns\n  • ML models classify packets\n  • Combined results determine attack type",
        "Alert Generation: If malicious activity detected, alerts are generated with severity levels",
        "Automated Response: Malicious IPs are automatically blocked via firewall",
        "User Notification: Real-time alerts sent to dashboard via WebSocket",
        "Data Storage: All packets and detections stored in MongoDB for audit trail"
    ])
    
    # Slide 13: Demo Scenarios
    slide = add_bullet_slide(prs, "Demo Scenarios", [
        "Scenario 1: Port Scan Detection\n  • Attacker performs port scan using nmap\n  • System detects probe attack pattern\n  • Alert appears in dashboard\n  • Source IP automatically blocked",
        "Scenario 2: DoS Attack Detection\n  • Attacker launches SYN flood attack\n  • System detects DoS pattern\n  • High severity alert generated\n  • IP blocked immediately",
        "Scenario 3: Brute Force Detection\n  • Multiple failed login attempts detected\n  • System identifies brute force pattern\n  • Alert with source IP displayed\n  • IP added to blocked list"
    ])
    
    # Slide 14: Results & Performance
    slide = add_bullet_slide(prs, "Results & Performance", [
        "Detection Accuracy: Binary ML model achieves 100% accuracy on test data",
        "Processing Speed: Real-time packet processing with <100ms detection latency",
        "System Reliability: Stable operation with rule-based detection as primary method",
        "User Experience: Intuitive dashboard with real-time updates and responsive design",
        "Scalability: Architecture supports horizontal scaling for high-traffic networks"
    ])
    
    # Slide 15: Challenges & Solutions
    slide = add_bullet_slide(prs, "Challenges & Solutions", [
        "Challenge 1: High Packet Volume\n  • Problem: Large number of packets causing system slowdown\n  • Solution: Efficient packet filtering and early exit for normal traffic",
        "Challenge 2: ML Model Accuracy\n  • Problem: Multiclass model struggled with attack classification\n  • Solution: Hybrid approach using rule-based detection with ML validation",
        "Challenge 3: Real-Time UI Updates\n  • Problem: Dashboard freezing with high packet volume\n  • Solution: Throttling, WebSocket optimization, and efficient rendering",
        "Challenge 4: False Positives\n  • Problem: Too many alerts overwhelming users\n  • Solution: Confidence scoring, severity levels, and duplicate detection"
    ])
    
    # Slide 16: Future Enhancements
    slide = add_bullet_slide(prs, "Future Enhancements", [
        "Deep Learning Models: Implement neural networks for better anomaly detection",
        "Geographic IP Tracking: Visualize attack sources on world map",
        "Attack Correlation: Identify related attacks and attack campaigns",
        "Mobile Application: Mobile app for remote monitoring",
        "SIEM Integration: Integration with Security Information and Event Management systems",
        "Advanced Analytics: Machine learning-based attack pattern prediction"
    ])
    
    # Slide 17: Conclusion
    slide = add_bullet_slide(prs, "Conclusion", [
        "Real-time network threat detection with high accuracy",
        "Automated response capabilities for immediate threat mitigation",
        "Intuitive web interface for security professionals",
        "Scalable architecture for enterprise deployment",
        "Comprehensive attack coverage for multiple threat types",
        "\nThe system successfully combines rule-based detection with machine learning to create a robust, reliable, and user-friendly network security solution."
    ])
    
    # Slide 18: Questions
    slide = add_text_slide(prs, "Questions?", "Thank you for your attention!")
    
    # Save presentation
    prs.save('IDS_Presentation.pptx')
    print("PowerPoint presentation created successfully: IDS_Presentation.pptx")

def add_text_slide(prs, title_text, content_text):
    """Add a slide with title and text content"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    title = slide.shapes.title
    title.text = title_text
    
    # Add text box for content
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    # Split content by newlines and add paragraphs
    paragraphs = content_text.split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = para_text
        p.font.size = Pt(12)
        p.space_after = Pt(6)
    
    return slide

def add_bullet_slide(prs, title_text, bullet_points):
    """Add a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    title = slide.shapes.title
    title.text = title_text
    
    # Add text box for bullet points
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.5)
    
    # Add bullet points
    for i, bullet_text in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(11)
        p.space_after = Pt(8)
        
        # Handle nested bullets (lines starting with "  •")
        lines = bullet_text.split('\n')
        if len(lines) > 1:
            for line in lines[1:]:
                if line.strip().startswith('•') or line.strip().startswith('-'):
                    nested_p = text_frame.add_paragraph()
                    nested_p.text = line.strip()
                    nested_p.level = 1
                    nested_p.font.size = Pt(10)
                    nested_p.space_after = Pt(4)
    
    return slide

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Error: python-pptx library not found.")
        print("Please install it using: pip install python-pptx")
    except Exception as e:
        print(f"Error creating presentation: {e}")




